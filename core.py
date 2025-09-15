import pandas as pd
import fitz  # PyMuPDF
import os
import re
import json
import logging
from pathlib import Path
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import CatchExcelImageTool


class ExcelToPDFProcessor:
    """Excel转PDF表单处理核心类"""
    
    def __init__(self):
        self.excel_path = ""
        self.pdf_template_path = ""
        self.output_folder = ""
        self.sheet_name = None
        self.title_row = 3
        self.start_row = 4
        self.filename_column = None  # 新增：文件名指定列
        self.field_mapping = {}
        self.col_separator = " "
        self.flatten_form = True
        self.output_png = False
        self.output_ppt = False

        # 字体相关配置
        self.font_base_path = "resources/fonts"  # 字体库基础路径
        self.default_font = "calibri"  # 默认字体（不含扩展名）
        self.chinese_font = "simhei"  # 中文字体（不含扩展名）
        self.default_fonts = {}  # 默认字体字典 {font_name: font_path}
        self.chinese_fonts = {}  # 中文字体字典 {font_name: font_path}

        # GUI日志回调函数
        self.gui_log_callback = None

        # 设置日志记录
        self.setup_logging()
        
        # 初始化字体库
        self.load_available_fonts()

    def _parse_widget_fontsize(self, widget, default=22):
        """
        从PDF widget中解析字体大小。支持数值或带'pt'的字符串（如'30 pt'）。
        将尝试以下属性名：field_fontsize / text_fontsize / fontsize。
        """
        candidates = []
        for attr in ("field_fontsize", "text_fontsize", "fontsize"):
            if hasattr(widget, attr):
                candidates.append(getattr(widget, attr))
        for val in candidates:
            if val is None:
                continue
            try:
                if isinstance(val, (int, float)) and val > 0:
                    return float(val)
                if isinstance(val, str):
                    s = val.strip().lower()
                    if s.endswith("pt"):
                        s = s[:-2].strip()
                    num = float(s)
                    if num > 0:
                        return num
            except:
                continue
        return float(default)
        
    def excel_col_letter_to_index(self, col):
        """将Excel列字母如'A'转为列号索引，从0开始"""
        result = 0
        for c in col:
            if 'A' <= c <= 'Z':
                result = result * 26 + (ord(c) - ord('A') + 1)
            elif 'a' <= c <= 'z':
                result = result * 26 + (ord(c) - ord('a') + 1)
            else:
                return None
        return result - 1
    
    def setup_logging(self):
        """设置日志记录"""
        # 创建日志记录器
        self.logger = logging.getLogger('ExcelToPDFProcessor')
        self.logger.setLevel(logging.DEBUG)
        
        # 如果已经有处理器，先清除
        if self.logger.handlers:
            self.logger.handlers.clear()
        
        # 创建文件处理器
        log_file = 'app.log'
        file_handler = logging.FileHandler(log_file, encoding='utf-8')
        file_handler.setLevel(logging.DEBUG)
        
        # 创建格式化器
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
            datefmt='%Y-%m-%d %H:%M:%S'
        )
        file_handler.setFormatter(formatter)
        
        # 添加处理器到记录器
        self.logger.addHandler(file_handler)
        
        self.logger.info("日志系统初始化完成")
    
    def load_available_fonts(self):
        """加载可用字体库"""
        self.default_fonts = {}  # 默认字体字典
        self.chinese_fonts = {}  # 中文字体字典
        
        # 加载默认字体
        default_path = os.path.join(self.font_base_path, "default")
        if os.path.exists(default_path):
            self._load_fonts_from_dir(default_path, self.default_fonts)
        
        # 加载中文字体
        zh_path = os.path.join(self.font_base_path, "zh")
        if os.path.exists(zh_path):
            self._load_fonts_from_dir(zh_path, self.chinese_fonts)
        
        total_fonts = len(self.default_fonts) + len(self.chinese_fonts)
        self.logger.info(f"加载字体库完成，默认字体: {len(self.default_fonts)}个，中文字体: {len(self.chinese_fonts)}个")
    
    def _load_fonts_from_dir(self, dir_path, font_dict):
        """从指定目录加载字体文件"""
        try:
            for file in os.listdir(dir_path):
                if file.lower().endswith(('.ttf', '.otf')):
                    font_path = os.path.join(dir_path, file)
                    font_name = os.path.splitext(file)[0]  # 直接使用文件名，不加前缀
                    font_dict[font_name] = font_path
        except Exception as e:
            self.logger.error(f"加载字体目录 {dir_path} 失败: {e}")
    
    def get_default_fonts(self):
        """获取默认字体列表"""
        return list(self.default_fonts.keys())
    
    def get_chinese_fonts(self):
        """获取中文字体列表"""
        return list(self.chinese_fonts.keys())
    
    def get_font_path(self, font_name, is_chinese=False):
        """根据字体名称获取字体文件路径"""
        if is_chinese:
            return self.chinese_fonts.get(font_name, "")
        else:
            return self.default_fonts.get(font_name, "")
    
    def has_chinese_characters(self, text):
        """检测文本中是否包含中文字符"""
        if not text:
            return False
        
        # 中文字符的Unicode范围
        chinese_pattern = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff]')
        return bool(chinese_pattern.search(text))
    
    def get_appropriate_font_path(self, text, default_font_name, chinese_font_name):
        """根据文本内容选择合适的字体路径"""
        has_chinese = self.has_chinese_characters(text)
        
        self.logger.debug(f"字体路径选择: 文本='{text[:20]}{'...' if len(text) > 20 else ''}', 包含中文={has_chinese}")
        
        if has_chinese and chinese_font_name:
            font_path = self.get_font_path(chinese_font_name, is_chinese=True)
            self.logger.debug(f"  尝试中文字体 '{chinese_font_name}': 路径={font_path}")
            if font_path and os.path.exists(font_path):
                self.logger.debug(f"  ✓ 使用中文字体: {font_path}")
                return font_path
            else:
                self.logger.debug(f"  ✗ 中文字体不可用")
        
        # 使用默认字体
        if default_font_name:
            font_path = self.get_font_path(default_font_name, is_chinese=False)
            self.logger.debug(f"  尝试默认字体 '{default_font_name}': 路径={font_path}")
            if font_path and os.path.exists(font_path):
                self.logger.debug(f"  ✓ 使用默认字体: {font_path}")
                return font_path
            else:
                self.logger.debug(f"  ✗ 默认字体不可用")
        
        # 如果都没有，返回None使用系统默认字体
        self.logger.debug(f"  → 回退到系统默认字体")
        return None
    
    def set_gui_log_callback(self, callback):
        """设置GUI日志回调函数"""
        self.gui_log_callback = callback
    
    def log_to_gui(self, operation, level="info", message=""):
        """发送日志到GUI（如果回调函数存在）"""
        if self.gui_log_callback:
            try:
                self.gui_log_callback(operation, level, message)
            except Exception as e:
                self.logger.error(f"GUI日志回调失败: {e}")

    def is_excel_col_pattern(self, val):
        """判断字符串是不是Excel列格式（如'A'或'A,B'）"""
        return re.fullmatch(r'([A-Za-z]+)(\s*,\s*[A-Za-z]+)*', val.strip()) is not None
    
    def extract_dispimg_id(self, cell_value):
        """从DISPIMG函数中提取图片ID"""
        if not cell_value or not isinstance(cell_value, str):
            return None
        
        # 匹配DISPIMG函数格式：=DISPIMG("ID_xxx",1) 或 =_xlfn.DISPIMG("ID_xxx",1)
        pattern = r'=(?:_xlfn\.)?DISPIMG\("([^"]+)"'
        match = re.search(pattern, str(cell_value))
        if match:
            return match.group(1)
        return None

    def get_pdf_form_keys(self, pdf_path):
        """获取PDF表单字段列表"""
        self.logger.info(f"开始获取PDF表单字段: {pdf_path}")
        
        if not os.path.exists(pdf_path):
            error_msg = f"PDF文件不存在: {pdf_path}"
            self.logger.error(error_msg)
            raise FileNotFoundError(error_msg)
            
        doc = fitz.open(pdf_path)
        fields = []
        
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                widgets = page.widgets()
                
                for widget in widgets:
                    field_name = widget.field_name
                    if field_name and field_name not in fields:
                        fields.append(field_name)
            
            self.logger.info(f"成功获取到 {len(fields)} 个PDF表单字段: {fields}")
        except Exception as e:
            error_msg = f"检查表单字段时出错: {e}"
            self.logger.error(error_msg)
            raise Exception(error_msg)
        finally:
            doc.close()
        
        return fields
    
    def fill_pdf_image_field(self, page, widget, image_path):
        """在PDF表单域中填充图片"""
        try:
            # 获取字段的矩形区域
            field_rect = widget.rect
            
            # 插入图片到字段区域
            page.insert_image(field_rect, filename=image_path, overlay=True)
            
            self.logger.debug(f"成功在字段 '{widget.field_name}' 中插入图片: {image_path}")
            return True
        except Exception as e:
            self.logger.error(f"在字段 '{widget.field_name}' 中插入图片失败: {e}")
            return False

    def fill_form_field_with_font(self, widget, value):
        """使用自定义字体填充表单字段（非扁平化模式）"""
        try:
            # 获取合适的字体路径
            font_path = self.get_appropriate_font_path(value, self.default_font, self.chinese_font)
            
            # 设置字段值
            widget.field_value = value
            
            # 如果有自定义字体，尝试设置字体
            if font_path and os.path.exists(font_path):
                try:
                    # 尝试设置字体（这个功能在某些PDF中可能不完全支持）
                    # PyMuPDF的表单字段字体设置有限，但我们仍然尝试
                    widget.update()
                    self.logger.debug(f"字段 '{widget.field_name}' 使用字体: {font_path}")
                except Exception as font_error:
                    # 如果字体设置失败，仍然使用默认字体填充
                    self.logger.debug(f"字段 '{widget.field_name}' 字体设置失败，使用默认字体: {font_error}")
            else:
                # 使用默认字体
                widget.update()
                self.logger.debug(f"字段 '{widget.field_name}' 使用默认字体")
            
            return True
            
        except Exception as e:
            self.logger.error(f"填充表单字段 '{widget.field_name}' 失败: {e}")
            return False

    def fill_pdf_form(self, input_pdf_path, output_pdf_path, data_dict, image_dict=None, flatten_form=False):
        """使用PyMuPDF填充PDF表单"""
        self.logger.debug(f"开始填充PDF表单: {output_pdf_path}")
        self.logger.debug(f"填充数据: {data_dict}")
        
        try:
            doc = fitz.open(input_pdf_path)
            
            filled_count = 0
            for page_num in range(len(doc)):
                page = doc[page_num]
                widgets = page.widgets()
                
                for widget in widgets:
                    field_name = widget.field_name
                    
                    # 处理图片字段
                    if image_dict and field_name in image_dict:
                        image_path = image_dict[field_name]
                        if image_path and os.path.exists(image_path):
                            if self.fill_pdf_image_field(page, widget, image_path):
                                filled_count += 1
                        continue
                    
                    # 处理普通文本字段
                    if field_name in data_dict:
                        value = str(data_dict[field_name])
                        if value:  # 只有非空值才填充
                            try:
                                if flatten_form:
                                    # 扁平化表单时只记录，实际填充在后面的扁平化过程中进行
                                    filled_count += 1
                                    self.logger.debug(f"准备处理字段 '{field_name}': 值='{value}'")
                                else:
                                    # 非扁平化表单：直接填充到表单字段中，支持字体选择
                                    success = self.fill_form_field_with_font(widget, value)
                                    if success:
                                        filled_count += 1
                                        self.logger.debug(f"成功填充字段 '{field_name}': 值='{value}'")
                                    else:
                                        self.logger.warning(f"填充字段 '{field_name}' 失败")
                            except Exception as e:
                                self.logger.error(f"填充字段 '{field_name}' 时出错: {e}")
            
            # 确保输出目录存在
            os.makedirs(os.path.dirname(output_pdf_path), exist_ok=True)
            
            # 保存文档
            if flatten_form:
                # 扁平化表单：优先使用 convert_to_pdf；失败时回退为“栅格化扁平化”（将整页渲染为图片）
                self.logger.debug("正在扁平化表单...")
                try:
                    self.flatten_form_with_textbox(doc, output_pdf_path, data_dict, image_dict)
                    self.logger.debug("表单扁平化完成")
                except Exception as e:
                    self.logger.error(f"textbox 扁平化失败，启用栅格化备用方案: {e}")
                    self.log_to_gui("扁平化处理", "warning", f"textbox扁平化失败，启用栅格化备用方案: {str(e)[:50]}...")
                    # 注意：该方式会生成不可编辑、不可检索的图像PDF，但能彻底避免字体替代错误
                    self.rasterize_flatten_doc(doc, output_pdf_path, dpi=200)
                    self.log_to_gui("扁平化处理", "info", "栅格化扁平化完成（图像PDF，不可编辑）")
            else:
                doc.save(output_pdf_path, incremental=False, encryption=fitz.PDF_ENCRYPT_NONE)
            
            doc.close()
            success_msg = f"成功填充了 {filled_count} 个字段"
            self.logger.info(f"PDF填充完成: {output_pdf_path}, {success_msg}")
            return True, success_msg
            
        except Exception as e:
            error_msg = f"PDF填充失败: {e}"
            self.logger.error(f"填充PDF失败 {output_pdf_path}: {error_msg}")
            return False, error_msg

    def process_excel_to_pdf(self, progress_callback=None):
        """主要处理函数：将Excel数据填充到PDF表单"""
        self.logger.info("开始处理Excel转PDF任务")
        self.logger.info(f"Excel文件: {self.excel_path}")
        self.logger.info(f"PDF模板: {self.pdf_template_path}")
        self.logger.info(f"输出目录: {self.output_folder}")
        self.logger.info(f"配置参数: 表头行={self.title_row}, 数据起始行={self.start_row}, 扁平化={self.flatten_form}")
        
        try:
            # 验证输入
            if not os.path.exists(self.excel_path):
                error_msg = f"Excel文件不存在: {self.excel_path}"
                self.logger.error(error_msg)
                raise FileNotFoundError(error_msg)
            if not os.path.exists(self.pdf_template_path):
                error_msg = f"PDF模板不存在: {self.pdf_template_path}"
                self.logger.error(error_msg)
                raise FileNotFoundError(error_msg)
            if not self.field_mapping:
                error_msg = "字段映射不能为空"
                self.logger.error(error_msg)
                raise ValueError(error_msg)
            
            # 创建输出目录
            os.makedirs(self.output_folder, exist_ok=True)
            self.logger.info(f"输出目录已创建: {self.output_folder}")
            
            # 读取Excel数据
            # 处理sheet_name为空的情况
            sheet_to_read = self.sheet_name if self.sheet_name and self.sheet_name.strip() else 0
            self.logger.info(f"读取Excel sheet: {sheet_to_read if sheet_to_read != 0 else '第一个sheet'}")
            
            df = pd.read_excel(self.excel_path, sheet_name=sheet_to_read, header=None)
            
            # 如果返回的是字典（多个sheet），取第一个
            if isinstance(df, dict):
                sheet_names = list(df.keys())
                df = df[sheet_names[0]]
                self.logger.info(f"读取到多个sheet，使用第一个: {sheet_names[0]}")
            
            title_row_idx = self.title_row - 1
            data_start_idx = self.start_row - 1
            self.logger.info(f"Excel数据读取完成，总行数: {len(df)}")
            
            # 获取PDF表单字段
            pdf_keys = self.get_pdf_form_keys(self.pdf_template_path)
            
            # 检查字段映射
            missing_fields = []
            for pdf_key in self.field_mapping.keys():
                if pdf_key not in pdf_keys:
                    missing_fields.append(pdf_key)
            
            if missing_fields:
                warning_msg = f"警告：以下字段在PDF中不存在: {missing_fields}"
                self.logger.warning(warning_msg)
                print(warning_msg)
            
            # 处理每一行数据
            total_rows = len(df.iloc[data_start_idx:])
            success_count = 0
            error_messages = []
            generated_pdf_paths = []  # 存储成功生成的PDF文件路径
            self.logger.info(f"开始处理数据，共 {total_rows} 行")
            self.log_to_gui("开始数据处理", "info", f"共 {total_rows} 行数据待处理")
            
            for idx, row in df.iloc[data_start_idx:].iterrows():
                row_num = idx + 1
                actual_row_num = idx - data_start_idx + 1
                
                try:
                    self.logger.info(f"开始处理第 {row_num} 行数据 (数据行 {actual_row_num}/{total_rows})")
                    self.log_to_gui(f"处理第{row_num}行", "info", f"开始处理 (数据行 {actual_row_num}/{total_rows})")
                    
                    # 构建数据字典和图片字典
                    data = {}
                    image_data = {}
                    field_details = []  # 记录字段映射详情
                    
                    for pdf_key, map_config in self.field_mapping.items():
                        if isinstance(map_config, dict):
                            is_excel_col = map_config.get("is_excel_col", True)
                            is_excel_image = map_config.get("is_excel_image", False)
                            map_val = str(map_config.get("val", ""))
                            # 只有Excel列才需要strip，自定义值保留原始格式（包括空格）
                            if is_excel_col or is_excel_image:
                                map_val = map_val.strip()
                        else:
                            # 向后兼容旧格式
                            map_val = str(map_config)
                            # 判断是否为Excel列格式，如果是才strip
                            if self.is_excel_col_pattern(map_val.strip()):
                                map_val = map_val.strip()
                                is_excel_col = True
                                is_excel_image = False
                            else:
                                is_excel_col = False
                                is_excel_image = False
                        
                        if is_excel_image:
                            # 处理Excel图片列
                            col_idx = self.excel_col_letter_to_index(map_val.upper())
                            if col_idx is not None and col_idx < len(row):
                                cell_val = row[col_idx]
                                # 对于图片字段，即使单元格为空也要尝试提取浮动图片
                                should_extract_image = True
                                
                                if should_extract_image:
                                    # 创建临时图片目录
                                    temp_image_dir = os.path.join(self.output_folder, "temp_images")
                                    os.makedirs(temp_image_dir, exist_ok=True)
                                    
                                    # 首先尝试提取DISPIMG中的图片ID（仅当单元格不为空时）
                                    image_id = None
                                    image_path = None
                                    
                                    if not pd.isna(cell_val) and cell_val:
                                        image_id = self.extract_dispimg_id(str(cell_val))
                                    
                                    if image_id:
                                        # DISPIMG格式图片
                                        image_path = CatchExcelImageTool.extract_image_by_id(
                                            self.excel_path, image_id, temp_image_dir
                                        )
                                        if image_path:
                                            image_data[pdf_key] = image_path
                                            field_details.append(f"{pdf_key}={map_val}(DISPIMG-ID:{image_id})→图片已提取")
                                        else:
                                            field_details.append(f"{pdf_key}={map_val}(DISPIMG-ID:{image_id})→图片提取失败")
                                    
                                    # 如果没有DISPIMG或DISPIMG提取失败，尝试提取浮动图片
                                    if not image_path:
                                        # 尝试提取浮动图片
                                        # 将列索引转换为单元格地址
                                        from openpyxl.utils import get_column_letter
                                        cell_address = f"{get_column_letter(col_idx + 1)}{row_num}"
                                        
                                        # 获取工作表名称
                                        sheet_to_use = self.sheet_name if self.sheet_name and self.sheet_name.strip() else None
                                        if not sheet_to_use:
                                            # 如果没有指定工作表名，尝试获取第一个工作表名
                                            try:
                                                import openpyxl
                                                wb_temp = openpyxl.load_workbook(self.excel_path, read_only=True)
                                                sheet_to_use = wb_temp.sheetnames[0]
                                                wb_temp.close()
                                            except:
                                                sheet_to_use = "Sheet1"  # 默认工作表名
                                        
                                        self.logger.debug(f"尝试从单元格 {cell_address} (工作表: {sheet_to_use}) 提取浮动图片")
                                        
                                        # 先检查是否有浮动图片存在
                                        floating_images = CatchExcelImageTool._extract_floating_images(self.excel_path)
                                        self.logger.debug(f"Excel文件中发现的浮动图片: {floating_images}")
                                        
                                        image_path = CatchExcelImageTool.extract_image_from_cell(
                                            self.excel_path, sheet_to_use, cell_address, temp_image_dir
                                        )
                                        
                                        if image_path:
                                            image_data[pdf_key] = image_path
                                            field_details.append(f"{pdf_key}={map_val}(浮动图片-{cell_address})→图片已提取")
                                            self.logger.info(f"成功提取浮动图片: {image_path}")
                                        else:
                                            # 提供详细的诊断信息，但不使用备用方案
                                            if floating_images:
                                                self.logger.warning(f"发现{len(floating_images)}个浮动图片但位置与单元格{cell_address}不匹配")
                                                field_details.append(f"{pdf_key}={map_val}(单元格{cell_address})→发现{len(floating_images)}个浮动图片但位置不匹配")
                                            else:
                                                field_details.append(f"{pdf_key}={map_val}(单元格{cell_address})→未发现任何图片")
                                                self.logger.warning(f"单元格 {cell_address} 未发现任何图片")
                                # 如果最终没有提取到任何图片，记录相应信息
                                if not image_path:
                                    if pd.isna(cell_val) or not cell_val:
                                        field_details.append(f"{pdf_key}={map_val}→单元格为空且未找到浮动图片")
                                    else:
                                        field_details.append(f"{pdf_key}={map_val}→单元格有值但未找到图片")
                            else:
                                field_details.append(f"{pdf_key}={map_val}→列索引无效")
                        elif is_excel_col:
                            # 处理Excel列值
                            col_list = [col.strip().upper() for col in map_val.split(",")]
                            cell_values = []
                            for col in col_list:
                                col_idx = self.excel_col_letter_to_index(col)
                                if col_idx is not None and col_idx < len(row):
                                    val = row[col_idx]
                                    cell_val = "" if pd.isna(val) else str(val)
                                    cell_values.append(cell_val)
                                else:
                                    cell_values.append("")
                            final_value = self.col_separator.join(cell_values)
                            data[pdf_key] = final_value
                            field_details.append(f"{pdf_key}={map_val}→'{final_value}'")
                        else:
                            # 直接使用自定义值，保留空格
                            data[pdf_key] = map_val
                            field_details.append(f"{pdf_key}=自定义值→'{map_val}'")
                    
                    self.logger.debug(f"第 {row_num} 行字段映射详情: {'; '.join(field_details)}")
                    
                    # 生成输出文件名
                    if self.filename_column and self.filename_column.strip():
                        # 使用指定列的数据作为文件名
                        try:
                            filename_col_idx = self.excel_col_letter_to_index(self.filename_column.strip())
                            if filename_col_idx is not None and filename_col_idx < len(row):
                                filename_value = str(row.iloc[filename_col_idx]).strip()
                                # 清理文件名中的非法字符
                                filename_value = re.sub(r'[<>:"/\\|?*]', '_', filename_value)
                                if filename_value:
                                    output_filename = f"{filename_value}.pdf"
                                else:
                                    output_filename = f"{actual_row_num}.pdf"  # 去掉filled_前缀，只保留数字
                            else:
                                output_filename = f"{actual_row_num}.pdf"  # 去掉filled_前缀，只保留数字
                        except:
                            output_filename = f"{actual_row_num}.pdf"  # 去掉filled_前缀，只保留数字
                    else:
                        # 使用数字编号，去掉filled_前缀
                        output_filename = f"{actual_row_num}.pdf"
                    
                    output_pdf_path = os.path.join(self.output_folder, output_filename)
                    
                    # 填充PDF表单
                    success, message = self.fill_pdf_form(
                        self.pdf_template_path, output_pdf_path, data, image_data, self.flatten_form
                    )
                    
                    if success:
                        success_count += 1
                        generated_pdf_paths.append(output_pdf_path)  # 记录成功生成的PDF路径
                        self.logger.info(f"第 {row_num} 行处理成功 → {os.path.basename(output_pdf_path)}")
                        self.log_to_gui(f"第{row_num}行", "success", f"处理成功 → {os.path.basename(output_pdf_path)}")
                    else:
                        error_msg = f"第{row_num}行: {message}"
                        error_messages.append(error_msg)
                        self.logger.error(f"第 {row_num} 行处理失败: {message}")
                        self.log_to_gui(f"第{row_num}行", "error", f"处理失败: {message}")
                    
                    # 更新进度
                    if progress_callback:
                        progress = actual_row_num / total_rows * 100
                        progress_callback(progress, f"处理第 {row_num} 行")
                        
                except Exception as e:
                    error_msg = f"第{row_num}行处理失败: {e}"
                    error_messages.append(error_msg)
                    self.logger.error(f"第 {row_num} 行处理异常: {str(e)}")
                    self.log_to_gui(f"第{row_num}行", "error", f"处理异常: {str(e)}")
            
            # 记录处理结果
            self.logger.info(f"处理完成 - 总行数: {total_rows}, 成功: {success_count}, 失败: {len(error_messages)}")
            if error_messages:
                self.logger.warning(f"处理过程中的错误: {error_messages}")
            
            # 发送汇总日志到GUI
            if len(error_messages) == 0:
                self.log_to_gui("处理完成", "success", f"全部 {total_rows} 行数据处理成功")
            else:
                self.log_to_gui("处理完成", "warning", f"总行数: {total_rows}, 成功: {success_count}, 失败: {len(error_messages)}")
                # 发送前几个错误详情到GUI
                for i, error_msg in enumerate(error_messages[:3]):
                    self.log_to_gui("错误详情", "error", error_msg)
                if len(error_messages) > 3:
                    self.log_to_gui("错误详情", "warning", f"还有 {len(error_messages) - 3} 个错误未显示")
            
            # 处理PNG和PPT转换
            png_paths = []
            ppt_path = None
            
            if generated_pdf_paths:
                # 转换为PNG图片
                if self.output_png:
                    self.log_to_gui("PNG转换", "info", "开始转换PDF为PNG图片...")
                    for pdf_path in generated_pdf_paths:
                        png_path = self.convert_pdf_to_png(pdf_path, self.output_folder)
                        if png_path:
                            png_paths.append(png_path)
                    
                    if png_paths:
                        self.log_to_gui("PNG转换", "success", f"成功转换 {len(png_paths)} 个PNG图片")
                    else:
                        self.log_to_gui("PNG转换", "error", "PNG转换失败")
                
                # 创建PPT文件
                if self.output_ppt:
                    self.log_to_gui("PPT创建", "info", "开始创建PPT文件...")
                    ppt_path = self.create_ppt_from_pdfs(generated_pdf_paths, self.output_folder)
                    
                    if ppt_path:
                        self.log_to_gui("PPT创建", "success", "PPT文件创建成功")
                    else:
                        self.log_to_gui("PPT创建", "error", "PPT文件创建失败")
            
            return {
                "success": True,
                "total_rows": total_rows,
                "success_count": success_count,
                "error_count": len(error_messages),
                "error_messages": error_messages,
                "generated_pdf_paths": generated_pdf_paths,
                "png_paths": png_paths,
                "ppt_path": ppt_path
            }
            
        except Exception as e:
            error_msg = str(e)
            self.logger.error(f"处理过程中发生严重错误: {error_msg}")
            return {
                "success": False,
                "error": error_msg
            }

    def save_preset(self, preset_path):
        """保存预设配置到JSON文件"""
        self.logger.info(f"开始保存预设配置到: {preset_path}")
        
        preset_data = {
            "excel_path": self.excel_path,
            "pdf_template_path": self.pdf_template_path,
            "output_folder": self.output_folder,
            "sheet_name": self.sheet_name,
            "title_row": self.title_row,
            "start_row": self.start_row,
            "filename_column": self.filename_column,  # 新增：文件名列
            "field_mapping": self.field_mapping,
            "col_separator": self.col_separator,
            "flatten_form": self.flatten_form,
            "output_png": self.output_png,
            "output_ppt": self.output_ppt,
            # 字体配置
            "font_base_path": self.font_base_path,  # 新增：字体库路径
            "default_font": self.default_font,
            "chinese_font": self.chinese_font
        }
        
        try:
            with open(preset_path, 'w', encoding='utf-8') as f:
                json.dump(preset_data, f, ensure_ascii=False, indent=2)
            success_msg = "预设保存成功"
            self.logger.info(f"预设配置保存成功: {preset_path}")
            return True, success_msg
        except Exception as e:
            error_msg = f"预设保存失败: {e}"
            self.logger.error(f"保存预设配置失败 {preset_path}: {error_msg}")
            return False, error_msg

    def load_preset(self, preset_path):
        """从JSON文件加载预设配置"""
        self.logger.info(f"开始加载预设配置: {preset_path}")
        
        try:
            with open(preset_path, 'r', encoding='utf-8') as f:
                preset_data = json.load(f)
            
            self.excel_path = preset_data.get("excel_path", "")
            self.pdf_template_path = preset_data.get("pdf_template_path", "")
            self.output_folder = preset_data.get("output_folder", "")
            self.sheet_name = preset_data.get("sheet_name", None)
            self.title_row = preset_data.get("title_row", 3)
            self.start_row = preset_data.get("start_row", 4)
            self.filename_column = preset_data.get("filename_column", None)  # 新增：文件名列
            self.field_mapping = preset_data.get("field_mapping", {})
            self.col_separator = preset_data.get("col_separator", " ")
            self.flatten_form = preset_data.get("flatten_form", False)
            self.output_png = preset_data.get("output_png", False)
            self.output_ppt = preset_data.get("output_ppt", False)
            
            # 加载字体配置
            self.font_base_path = preset_data.get("font_base_path", "resources/fonts")  # 新增：字体库路径
            self.default_font = preset_data.get("default_font", "calibri")
            self.chinese_font = preset_data.get("chinese_font", "simhei")
            
            # 重新加载字体库
            self.load_available_fonts()
            
            success_msg = "预设加载成功"
            self.logger.info(f"预设配置加载成功: {preset_path}")
            self.logger.debug(f"加载的配置: {preset_data}")
            return True, success_msg
        except Exception as e:
            error_msg = f"预设加载失败: {e}"
            self.logger.error(f"加载预设配置失败 {preset_path}: {error_msg}")
            return False, error_msg

    def reset_to_default(self):
        """恢复默认配置"""
        self.logger.info("开始恢复默认配置")
        
        self.excel_path = ""
        self.pdf_template_path = ""
        self.output_folder = str(Path.home() / "Desktop")
        self.sheet_name = None
        self.title_row = 3
        self.start_row = 4
        self.filename_column = None  # 新增：文件名列重置
        self.field_mapping = {}
        self.col_separator = " "
        self.flatten_form = True
        self.output_png = False
        self.output_ppt = False
        
        # 重置字体配置
        self.font_base_path = "resources/fonts"
        self.default_font = "calibri"
        self.chinese_font = "simhei"
        
        # 重新加载字体库
        self.load_available_fonts()
        
        self.logger.info("默认配置恢复完成")
        return "已恢复默认配置"

    def get_desktop_path(self):
        """获取桌面路径"""
        return str(Path.home() / "Desktop")
    


    def flatten_form_with_textbox(self, doc, output_pdf_path, data_dict, image_dict=None):
        """使用insert_text方式扁平化表单，参考用户示例代码的简单直接方法"""
        try:
            # 创建新文档，复制原文档的页面但不包含表单字段
            new_doc = fitz.open()
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # 创建新页面，复制原页面尺寸
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                
                # 复制原页面的内容（不包括表单字段）
                # 先渲染原页面为图片，然后插入到新页面
                pix = page.get_pixmap(alpha=False)
                new_page.insert_image(new_page.rect, pixmap=pix)
                
                # 字体嵌入字典，用于存储已嵌入的字体
                embedded_fonts = {}
                
                # 获取原页面的表单字段并在新页面上用insert_text重新绘制
                widgets = page.widgets()
                for widget in widgets:
                    field_name = widget.field_name
                    
                    # 处理图片字段
                    if image_dict and field_name in image_dict:
                        image_path = image_dict[field_name]
                        if image_path and os.path.exists(image_path):
                            try:
                                # 获取字段的矩形区域
                                field_rect = widget.rect
                                # 在新页面插入图片
                                new_page.insert_image(field_rect, filename=image_path, overlay=True)
                                self.logger.debug(f"成功在字段 '{field_name}' 位置插入图片")
                            except Exception as e:
                                self.logger.error(f"在字段 '{field_name}' 位置插入图片失败: {e}")
                        continue
                    
                    # 处理文本字段
                    if field_name in data_dict:
                        value = str(data_dict[field_name])
                        if value:  # 只有非空值才绘制
                            # 获取字段的矩形区域
                            field_rect = widget.rect
                            
                            try:
                                # 根据表单域设置的字体大小作为初始字号（支持'xx pt'）
                                fontfile = self.get_appropriate_font_path(value, self.default_font, self.chinese_font)
                                
                                # 获取原始字体大小，直接使用不进行调整
                                font_size = self._parse_widget_fontsize(widget, default=22)
                                color = (0, 0, 0)  # 黑色
                                
                                # 详细记录用户字体选择和应用情况
                                has_chinese = self.has_chinese_characters(value)
                                selected_font_name = self.chinese_font if has_chinese else self.default_font
                                
                                self.logger.info(f"=== 字段 '{field_name}' 字体应用详情 ===")
                                self.logger.info(f"  填充内容: '{value}' (长度: {len(value)})")
                                self.logger.info(f"  包含中文: {has_chinese}")
                                self.logger.info(f"  用户选择的默认字体: {self.default_font}")
                                self.logger.info(f"  用户选择的中文字体: {self.chinese_font}")
                                self.logger.info(f"  本次应该使用的字体: {selected_font_name}")
                                self.logger.info(f"  字体文件路径: {fontfile if fontfile else '未找到字体文件，将使用系统默认字体'}")
                                self.logger.info(f"  表单域字体大小: {font_size}pt")
                                
                                # 嵌入字体并获得fontname
                                fontname = None
                                if fontfile and os.path.exists(fontfile):
                                    # 检查是否已经嵌入过这个字体
                                    if fontfile not in embedded_fonts:
                                        try:
                                            # 先把字体嵌入并得到fontname
                                            fontname = f"CustomFont_{len(embedded_fonts)}"  # 生成唯一字体名
                                            xref = new_page.insert_font(fontname=fontname, fontfile=fontfile)
                                            embedded_fonts[fontfile] = fontname
                                            self.logger.info(f"  ✓ 字体嵌入成功: {os.path.basename(fontfile)} -> fontname={fontname}")
                                        except Exception as font_error:
                                            self.logger.warning(f"  ✗ 字体嵌入失败: {font_error}")
                                            fontname = None
                                    else:
                                        fontname = embedded_fonts[fontfile]
                                        self.logger.info(f"  ✓ 使用已嵌入字体: {os.path.basename(fontfile)} -> fontname={fontname}")
                                else:
                                    self.logger.warning(f"  ✗ 字体文件不存在或路径无效，将使用系统默认字体")
                                
                                # 计算垂直居中的Y位置（参考用户示例代码）
                                text_height = font_size
                                center_y = (field_rect.y0 + field_rect.y1) / 2
                                insert_x = field_rect.x0
                                insert_y = center_y - text_height / 2 + text_height * 0.75
                                
                                # 使用insert_text直接插入文本（参考用户示例代码）
                                new_page.insert_text(
                                    fitz.Point(insert_x, insert_y),
                                    value,
                                    fontname=fontname,
                                    fontsize=font_size,
                                    color=color,
                                )
                                
                                self.logger.info(f"  ✓ 文本插入成功: 字号={font_size}pt, 位置=({insert_x:.1f}, {insert_y:.1f})")
                                self.logger.info(f"  字段矩形: {field_rect}")
                                self.logger.info(f"=== 字段 '{field_name}' 处理完成 ===")
                                
                            except Exception as e:
                                self.logger.error(f"在字段 '{field_name}' 位置插入文本失败: {e}")
                                # 备用方案：使用系统默认字体
                                try:
                                    font_size = self._parse_widget_fontsize(widget, default=14)
                                    text_height = font_size
                                    center_y = (field_rect.y0 + field_rect.y1) / 2
                                    insert_x = field_rect.x0
                                    insert_y = center_y - text_height / 2 + text_height * 0.75
                                    
                                    self.logger.info(f"  === 启动备用方案（系统默认字体） ===")
                                    
                                    new_page.insert_text(
                                        fitz.Point(insert_x, insert_y),
                                        value,
                                        fontname=None,  # 使用系统默认字体
                                        fontsize=font_size,
                                        color=(0, 0, 0),
                                    )
                                    
                                    self.logger.info(f"  ✓ 备用方案成功: 字号={font_size}pt")
                                    self.logger.info(f"=== 字段 '{field_name}' 备用方案完成 ===")
                                except Exception as e2:
                                    self.logger.error(f"  ✗ 备用方案也失败: {e2}")
                                    self.logger.error(f"=== 字段 '{field_name}' 处理失败 ===")
            
            # 保存新文档
            new_doc.save(output_pdf_path, deflate=True, clean=True)
            new_doc.close()
            self.logger.debug(f"insert_text扁平化完成: {output_pdf_path}")
            
        except Exception as e:
            self.logger.error(f"insert_text扁平化失败: {e}")
            raise

    def rasterize_flatten_doc(self, doc, output_pdf_path, dpi=200):
        """将PDF以指定DPI栅格化后重新封装为不可编辑的图像PDF，避免字体替代问题"""
        try:
            zoom = dpi / 72.0  # 72pt = 1英寸
            mat = fitz.Matrix(zoom, zoom)
            new_doc = fitz.open()
            for page in doc:
                # 渲染为位图
                pix = page.get_pixmap(matrix=mat, alpha=False)
                # 新建与原页面同尺寸的页面（以pt为单位）
                new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                # 将渲染得到的图片铺满页面
                new_page.insert_image(new_page.rect, pixmap=pix)
            # 保存图像PDF
            new_doc.save(output_pdf_path, deflate=True, clean=True)
            new_doc.close()
            self.logger.debug(f"栅格化扁平化完成: {output_pdf_path}")
        except Exception as e:
            self.logger.error(f"栅格化扁平化失败: {e}")
            raise
    
    def convert_pdf_to_png(self, pdf_path, output_folder):
        """将PDF转换为PNG图片"""
        try:
            pdf_name = Path(pdf_path).stem
            png_path = os.path.join(output_folder, f"{pdf_name}.png")
            
            # 打开PDF文件
            pdf_doc = fitz.open(pdf_path)
            
            # 获取第一页（假设PDF只有一页）
            page = pdf_doc[0]
            
            # 设置缩放比例以获得高质量图片
            zoom = 2.0  # 缩放比例
            mat = fitz.Matrix(zoom, zoom)
            
            # 渲染页面为图片
            pix = page.get_pixmap(matrix=mat)
            
            # 保存为PNG
            pix.save(png_path)
            
            pdf_doc.close()
            
            self.logger.info(f"PDF转PNG成功: {png_path}")
            self.log_to_gui("PDF转PNG", "info", f"成功转换: {pdf_name}.png")
            
            return png_path
            
        except Exception as e:
            error_msg = f"PDF转PNG失败: {str(e)}"
            self.logger.error(error_msg)
            self.log_to_gui("PDF转PNG", "error", error_msg)
            return None
    
    def create_ppt_from_pdfs(self, pdf_paths, output_folder):
        """将多个PDF文件合并为一个PPT文件"""
        try:
            # 临时增加PIL的图像大小限制
            from PIL import Image
            original_max_image_pixels = Image.MAX_IMAGE_PIXELS
            Image.MAX_IMAGE_PIXELS = None  # 临时移除限制
            
            # 创建PPT演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸为A4比例
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for pdf_path in pdf_paths:
                try:
                    # 打开PDF文件
                    pdf_doc = fitz.open(pdf_path)
                    
                    # 获取第一页
                    page = pdf_doc[0]
                    
                    # 获取页面尺寸
                    page_rect = page.rect
                    page_width = page_rect.width
                    page_height = page_rect.height
                    
                    # 计算安全的缩放比例，避免生成过大的图像
                    # 限制最大像素数为150M像素（约为PIL默认限制的85%）
                    max_pixels = 150_000_000
                    current_pixels = page_width * page_height
                    
                    if current_pixels > 0:
                        max_zoom = (max_pixels / current_pixels) ** 0.5
                        # 选择较小的缩放比例，但至少为0.5以保证质量
                        zoom = min(2.0, max(0.5, max_zoom))
                    else:
                        zoom = 1.0
                    
                    self.logger.info(f"PDF页面尺寸: {page_width:.1f}x{page_height:.1f}, 使用缩放比例: {zoom:.2f}")
                    
                    mat = fitz.Matrix(zoom, zoom)
                    
                    # 渲染页面为图片
                    pix = page.get_pixmap(matrix=mat)
                    
                    # 检查生成的图像尺寸
                    img_width = pix.width
                    img_height = pix.height
                    total_pixels = img_width * img_height
                    
                    self.logger.info(f"生成图像尺寸: {img_width}x{img_height}, 总像素数: {total_pixels:,}")
                    
                    # 将图片数据转换为PIL Image
                    img_data = pix.tobytes("png")
                    
                    # 创建临时图片文件
                    temp_img_path = os.path.join(output_folder, f"temp_{Path(pdf_path).stem}.png")
                    with open(temp_img_path, "wb") as f:
                        f.write(img_data)
                    
                    # 添加新幻灯片
                    slide_layout = prs.slide_layouts[6]  # 空白布局
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 添加图片到幻灯片
                    left = Inches(0.5)
                    top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(6.5)
                    
                    slide.shapes.add_picture(temp_img_path, left, top, width, height)
                    
                    # 删除临时图片文件
                    os.remove(temp_img_path)
                    
                    pdf_doc.close()
                    
                    self.logger.info(f"成功处理PDF文件: {Path(pdf_path).name}")
                    
                except Exception as e:
                    error_msg = f"处理PDF文件 {pdf_path} 时出错: {str(e)}"
                    self.logger.warning(error_msg)
                    self.log_to_gui("PPT处理", "warning", f"跳过文件 {Path(pdf_path).name}: {str(e)[:50]}...")
                    continue
            
            # 恢复PIL的原始限制
            Image.MAX_IMAGE_PIXELS = original_max_image_pixels
            
            # 保存PPT文件
            ppt_path = os.path.join(output_folder, "merged_pdfs.pptx")
            prs.save(ppt_path)
            
            self.logger.info(f"PPT创建成功: {ppt_path}")
            self.log_to_gui("创建PPT", "info", f"成功创建PPT: merged_pdfs.pptx")
            
            return ppt_path
            
        except Exception as e:
            # 确保恢复PIL的原始限制
            try:
                from PIL import Image
                Image.MAX_IMAGE_PIXELS = original_max_image_pixels
            except:
                pass
            
            error_msg = f"创建PPT失败: {str(e)}"
            self.logger.error(error_msg)
            self.log_to_gui("创建PPT", "error", error_msg)
            return None